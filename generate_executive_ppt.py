from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


OUTPUT_FILE = "Comite_Internacional_Conferencia_System.pptx"


def set_background(slide, color=(245, 248, 252)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def style_title(shape, size=34, color=(11, 42, 74)):
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*color)


def add_footer(slide, text="Confidential | Executive Committee"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.2), Inches(0.3))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(90, 90, 90)
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, title, bullets, subtitle=None):
    layout = prs.slide_layouts[1]  # Title and Content
    s = prs.slides.add_slide(layout)
    set_background(s)

    s.shapes.title.text = title
    style_title(s.shapes.title, size=32)

    content = s.placeholders[1]
    tf = content.text_frame
    tf.clear()

    if subtitle:
        p0 = tf.paragraphs[0]
        p0.text = subtitle
        p0.font.bold = True
        p0.font.size = Pt(20)
        p0.font.color.rgb = RGBColor(28, 67, 106)

    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(30, 30, 30)

    add_footer(s)


def add_title_slide(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(s, color=(234, 242, 250))

    s.shapes.title.text = title
    style_title(s.shapes.title, size=40, color=(9, 35, 64))

    st = s.placeholders[1]
    st.text = subtitle
    p = st.text_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(30, 74, 116)

    add_footer(s, text="Prepared for International Executive Committee")


def add_agenda():
    add_bullets(
        slide=None,
        title="Agenda",
        subtitle="Conferência System | Executive Overview",
        bullets=[
            "1) System Overview and Business Value",
            "2) Functional Scope and Technical Architecture",
            "3) Operational Workflow and Innovations",
            "4) Current Status, Risks and Future Roadmap",
            "5) Real Use Cases and Executive Highlights",
        ],
    )


def build_presentation():
    add_title_slide(
        "Conferência System",
        "Integrated Receiving, Fiscal, WMS and Shipping Platform",
    )

    add_agenda()

    add_bullets(
        slide=None,
        title="1. System Overview",
        bullets=[
            "Business objective: unify receiving, fiscal release, warehousing and shipping in one platform.",
            "Problem solved: manual controls, low traceability, rework and compliance exposure.",
            "Target users: Gate, Receiving Operators, Fiscal Team, Managers and Warehouse Team.",
            "Operating model: web access from tablets and desktops in local network.",
        ],
    )

    add_bullets(
        slide=None,
        title="2. Business Value",
        bullets=[
            "Operational reliability: end-to-end traceability for all critical events.",
            "Efficiency gains: faster conferences and reduced manual reconciliation.",
            "Estimated impact: 30% to 60% cycle-time reduction in conference and pendency handling.",
            "Strategic value: scalable digital foundation for governance and international standards.",
        ],
    )

    add_bullets(
        slide=None,
        title="3. Functionalities | Core Capabilities",
        bullets=[
            "Role-based access control with granular permissions per profile and user.",
            "XML import, pending queue, conference locks and heartbeat protection.",
            "Complete timeline of reversals, exclusions, approvals and launch actions.",
            "Checklist, evidence upload and SLA dashboards for management visibility.",
        ],
    )

    add_bullets(
        slide=None,
        title="3. Functionalities | Automation and AI-like Logic",
        bullets=[
            "Smart CFOP filtering to avoid false pendencies in mixed fiscal notes.",
            "Automatic quantity tolerance (±2%) for KG and MM units.",
            "Rule-based XML auditor with fiscal diagnostics and guided decision support.",
            "Idempotent WMS integration queue with retry/dead-letter processing model.",
        ],
    )

    add_bullets(
        slide=None,
        title="4. Technical Architecture",
        bullets=[
            "Frontend: server-rendered HTML templates with operational pages for each role.",
            "Backend: Python + Flask, validation via Marshmallow, services by module.",
            "Data layer: SQLAlchemy ORM over SQLite (migration-ready to PostgreSQL).",
            "Production serving: Waitress; module APIs for conference, fiscal, WMS and shipping.",
        ],
    )

    add_bullets(
        slide=None,
        title="4. Integrations",
        bullets=[
            "Consyste API for consultation, document retrieval and recipient manifestation.",
            "Google Sheets purchase-order consultation with Excel/cache fallback.",
            "Internal REST APIs for dashboards, governance and operational controls.",
            "Tablet-first browser operation for low-cost, high-accessibility deployment.",
        ],
    )

    add_bullets(
        slide=None,
        title="5. Workflow | How It Operates",
        bullets=[
            "Step 1: Import NF XML and normalize fiscal/product data.",
            "Step 2: Execute fiscal-operational audit and queue pending actions.",
            "Step 3: Blind conference validates counts, tolerance and divergence reasons.",
            "Step 4: Fiscal launch triggers manifestation and WMS integration event.",
            "Step 5: WMS endereces items, reconciles ERP x WMS and raises operational alerts.",
            "Step 6: Shipping performs blind validation, partial/total invoicing and controlled reversals.",
        ],
    )

    add_bullets(
        slide=None,
        title="6. Innovations / Differentials",
        bullets=[
            "Blind conference with lock-control and anti-concurrency heartbeat.",
            "Unified governance layer across fiscal, receiving, warehouse and shipping.",
            "Operationally pragmatic architecture: low cost now, scalable roadmap next.",
            "Audit-ready logs by design: every critical action is traceable.",
        ],
    )

    add_bullets(
        slide=None,
        title="7. Current Status",
        bullets=[
            "Deployment level: operational use with active modules in daily routines.",
            "Quality signal: automated test suite with 41 mapped scenarios.",
            "Stability: medium-high for current scale, with controlled known limits.",
            "Known limitations: SQLite scale ceiling and dependency on external API availability.",
        ],
    )

    add_bullets(
        slide=None,
        title="8. Future Improvements",
        bullets=[
            "Migrate to PostgreSQL for higher concurrency and enterprise scalability.",
            "Add dedicated async workers for queue processing and scheduling.",
            "Expand KPIs: lead time by stage, SKU accuracy, supplier quality trends.",
            "Add advanced security controls: SSO, MFA and stronger observability.",
        ],
    )

    add_bullets(
        slide=None,
        title="9. Real Use Cases",
        bullets=[
            "Inbound note with mixed CFOP: system filters non-conference lines automatically.",
            "Divergence event: operator logs evidence and physical destination instantly.",
            "Fiscal release: recipient manifestation logged and launch integrated to WMS.",
            "Warehouse reconciliation: ERP vs WMS mismatch detected and tracked to resolution.",
        ],
    )

    add_bullets(
        slide=None,
        title="10. Data for Presentation",
        bullets=[
            "Key phrases: End-to-end traceability | Fiscal compliance by design | Event-driven integration.",
            "Highlights: 41 automated tests | ±2% smart tolerance (KG/MM) | Multi-module coverage.",
            "Executive message: rapid ROI with low-cost deployment and clear scale roadmap.",
            "Positioning: operations-first platform with governance and audit readiness.",
        ],
    )

    add_bullets(
        slide=None,
        title="Appendix | Suggested KPI Dashboard",
        bullets=[
            "Conference cycle time (median and p95).",
            "Divergence rate by supplier and SKU.",
            "Fiscal launch success rate and manifestation outcomes.",
            "WMS reconciliation open items and aging buckets.",
            "Shipping first-pass accuracy and reversal index.",
        ],
    )

    add_bullets(
        slide=None,
        title="Closing",
        bullets=[
            "Conferência System consolidates operational control across critical logistics-fiscal processes.",
            "Current architecture delivers immediate value while preserving enterprise evolution options.",
            "Recommendation: approve next-phase scaling (database, queue workers, advanced KPIs).",
        ],
        subtitle="Thank you",
    )


if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_presentation()
    prs.save(OUTPUT_FILE)
    print(f"Presentation generated: {OUTPUT_FILE}")
